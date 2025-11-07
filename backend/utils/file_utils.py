import os
import shutil
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont
from io import BytesIO

OUTPUT_DIR = "output/"

def ensure_output_dir():
    """Create output folder if not exists."""
    if not os.path.exists(OUTPUT_DIR):
        os.makedirs(OUTPUT_DIR)
    return OUTPUT_DIR

def clear_output_dir():
    """Delete and recreate output folder (clean previous runs)."""
    if os.path.exists(OUTPUT_DIR):
        shutil.rmtree(OUTPUT_DIR)
    os.makedirs(OUTPUT_DIR)

def save_file(content, filename):
    """Save text or binary content to file."""
    ensure_output_dir()
    filepath = os.path.join(OUTPUT_DIR, filename)
    mode = 'wb' if isinstance(content, bytes) else 'w'
    with open(filepath, mode) as f:
        f.write(content)
    return filepath

def list_output_files():
    """List all generated files."""
    ensure_output_dir()
    return os.listdir(OUTPUT_DIR)

def convert_pptx_to_images(pptx_path):
    """
    Converts slides in a PowerPoint file to images.
    This implementation renders basic slide content to images.
    """
    ensure_output_dir()  # Ensure output directory exists
    prs = Presentation(pptx_path)
    slide_images = []

    # Use a better font and size for readability
    try:
        # Try to use a better font if available
        font_title = ImageFont.truetype("arial.ttf", 36)
        font_content = ImageFont.truetype("arial.ttf", 24)
    except:
        # Fallback to default font if specific font is not available
        font_title = ImageFont.load_default()
        font_content = ImageFont.load_default()

    for i, slide in enumerate(prs.slides):
        img_path = os.path.join(OUTPUT_DIR, f"slide_{i+1}.png")
        
        # Create an image with slide dimensions
        width, height = 1280, 720  # Standard HD resolution
        img = Image.new("RGB", (width, height), color=(255, 255, 255))  # White background
        draw = ImageDraw.Draw(img)
        
        # Try to add some basic text from the slide
        y_offset = 50
        
        # Add title if exists
        if slide.shapes.title and hasattr(slide.shapes.title, 'text'):
            title_text = slide.shapes.title.text.strip()
            if title_text:
                draw.text((50, y_offset), title_text, fill=(0, 0, 0), font=font_title)
                y_offset += 80  # Increased spacing for title
        
        # Add content from placeholders
        for shape in slide.shapes:
            # Check if shape has text frame
            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                # Get text content (ignore linter warning as this works at runtime)
                text_frame = shape.text_frame  # type: ignore
                if hasattr(text_frame, 'text'):
                    text = text_frame.text.strip()  # type: ignore
                    if text and shape != slide.shapes.title:  # Skip title as we already handled it
                        try:
                            # Draw other text content with better formatting
                            lines = text.split('\n')
                            for line in lines:
                                if line.strip():
                                    draw.text((70, y_offset), line, fill=(50, 50, 50), font=font_content)
                                    y_offset += 40  # Increased spacing for better readability
                        except:
                            # Fallback if there are font issues
                            draw.text((50, y_offset), text[:100], fill=(0, 0, 0), font=font_content)  # Limit text length
                            y_offset += 40
        
        img.save(img_path)
        slide_images.append(img_path)
        print(f"Saved slide image: {img_path}, exists: {os.path.exists(img_path)}")

    # Upload images to Cloudinary
    try:
        from services.cloud_service import upload_file
        cloud_images = []
        for img_path in slide_images:
            upload_result = upload_file(img_path, resource_type="image")
            if upload_result:
                cloud_images.append(upload_result["secure_url"])
            else:
                # Fallback to local path if cloud upload fails
                cloud_images.append(img_path)
        # Return both local paths (for video creation) and cloud URLs (for response)
        return {"local_paths": slide_images, "cloud_urls": cloud_images}
    except Exception as e:
        print(f"Cloudinary upload failed: {e}")
        # Fallback to local paths if cloud upload fails
        return {"local_paths": slide_images, "cloud_urls": slide_images}
