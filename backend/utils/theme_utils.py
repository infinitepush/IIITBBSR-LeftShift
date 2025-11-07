from pptx.dml.color import RGBColor

def get_theme_settings(theme_name):
    """Return theme style settings based on user choice."""
    themes = {
        "Minimalist": {
            "bg_color": RGBColor(255, 255, 255),  # white
            "title_color": RGBColor(0, 0, 0),      # black
            "font": "Arial"
        },
        "Chalkboard": {
            "bg_color": RGBColor(28, 28, 28),      # dark gray/blackboard
            "title_color": RGBColor(255, 255, 255),# white text
            "font": "Comic Sans MS"
        },
        "Corporate": {
            "bg_color": RGBColor(245, 245, 245),   # light gray
            "title_color": RGBColor(0, 51, 102),   # navy blue
            "font": "Calibri"
        }
    }
    return themes.get(theme_name, themes["Minimalist"])


def apply_theme(prs, theme_name):
    """
    Apply background color and text color to PowerPoint presentation.
    This will affect all slides added after applying the theme.
    """
    theme = get_theme_settings(theme_name)

    # PowerPoint theme adjustment is limited — we’ll apply basic coloring here.
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = theme["font"]
                    run.font.color.rgb = theme["title_color"]

    # Background color — new slides will follow this pattern
    prs.slide_master.background.fill.solid()
    prs.slide_master.background.fill.fore_color.rgb = theme["bg_color"]
