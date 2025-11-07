from pptx import Presentation
from pptx.util import Inches
from utils.theme_utils import apply_theme
import os

def generate_slides(slides_data, theme):
    # Ensure output directory exists
    output_dir = "output"
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
        
    prs = Presentation()
    apply_theme(prs, theme)

    for slide in slides_data:
        slide_layout = prs.slide_layouts[1]
        s = prs.slides.add_slide(slide_layout)
        
        # Set title
        if s.shapes.title and 'title' in slide:
            s.shapes.title.text = slide['title']
            
        # Set content
        if len(s.placeholders) > 1:
            body = s.placeholders[1]
            if body.has_text_frame:
                text_frame = body.text_frame  # type: ignore
                text_frame.clear()  # Clear any existing text
                
                for i, point in enumerate(slide['content']):
                    if i == 0:
                        text_frame.text = f"• {point}"
                    else:
                        p = text_frame.add_paragraph()
                        p.text = f"• {point}"

    path = "output/slides.pptx"
    prs.save(path)
    return path